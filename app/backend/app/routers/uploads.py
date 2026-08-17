import logging
import uuid
from pathlib import Path

import pymupdf
from fastapi import APIRouter, Depends, HTTPException, Request, UploadFile, status
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession
from starlette.responses import FileResponse

from app.core.config import get_settings
from app.db import get_db
from app.deps import ensure_space_access, get_current_user, is_vault_space
from app.models import Upload, User
from app.schemas.upload import UploadOut
from app.pdf_processing import AUTO_OCR_MAX_PDF_BYTES
from app.pdf_processing import enqueue_ocr as enqueue_pdf_ocr
from app.pdf_processing import extract_text as extract_pdf_text
from app.transcription import enqueue_transcription
from app.vision import enqueue_vision
from app import vision as vision_module
from app.note_recording import _MIN_RECORDING_BYTES, enqueue_recording

router = APIRouter(prefix="/api/uploads", tags=["uploads"])
logger = logging.getLogger(__name__)

# Произвольные файлы (ТЗ §9 — файл как item — придёт в Фазе 2; пока это
# просто вложение к заметке). Тип не ограничиваем: файл никогда не
# исполняется, только отдаётся обратно через авторизованный FileResponse.
# 300 МБ хватает на видео с телефона в разумном качестве (диск — 27 ГБ
# свободно, не узкое место); раньше было 25 МБ — заимствовано у лимита
# Whisper API для голосовых сообщений, для видео категорически мало.
MAX_UPLOAD_BYTES = 300 * 1024 * 1024
_STREAM_CHUNK_BYTES = 1024 * 1024

# Реальная жалоба: "широкий набор файлов" должен получить хоть какое-то
# превью, не только картинки/видео/PDF. ВИЗУАЛЬНЫЙ рендер офисных форматов
# (docx/xlsx/pptx — "как страница реально выглядит") потребовал бы тяжёлую
# новую зависимость (LibreOffice headless — не просто конвертер, а весь
# офисный движок целиком, 300-500+ МБ RAM на конвертацию, а бюджет памяти
# уже расписан под ноль, см. CLAUDE.md) — сознательно не делаем. Но ТЕКСТ
# из них достать дёшево и без единого тяжёлого бинарника: python-docx/
# openpyxl/python-pptx — чистый Python (плюс lxml, уже лёгкая и обычная
# зависимость), тот же по духу приём, что и текстовый слой PDF (PyMuPDF)
# или чтение .txt/.md ниже — просто для другого формата хранения текста.
_TEXT_PREVIEW_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".log", ".yaml", ".yml", ".xml",
    ".py", ".js", ".ts", ".tsx", ".jsx", ".css", ".html", ".sh", ".ini", ".conf",
}
_OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
_MAX_TEXT_PREVIEW_BYTES = 200 * 1024
_MAX_TEXT_PREVIEW_CHARS = 20000
_MAX_THUMBNAIL_DPI = 100
_MAX_OFFICE_PREVIEW_ROWS = 300  # xlsx: не сканируем огромные таблицы целиком ради превью
_MAX_OFFICE_PREVIEW_BYTES = 20 * 1024 * 1024  # docx/xlsx/pptx крупнее — превью пропускаем, не парсим целиком


def _upload_path(upload_id: uuid.UUID) -> Path:
    return Path(get_settings().upload_dir) / str(upload_id)


def _thumbnail_path(upload_id: uuid.UUID) -> Path:
    return Path(get_settings().upload_dir) / f"{upload_id}.thumb.png"


def _staged_upload_path(upload_id: uuid.UUID) -> Path:
    # Смена пароля сейфа (routers/spaces.py::rotate_vault_password) — файл
    # перешифровывается сюда, а не поверх оригинала: реальный файл должен
    # остаться читаемым СТАРЫМ ключом, пока вся операция (все заметки +
    # сама соль/verifier спейса) не подтверждена одной транзакцией.
    # os.replace() поверх оригинала происходит только там, ПОСЛЕ commit.
    return Path(get_settings().upload_dir) / f"{upload_id}.new"


def _extract_docx_text(path: Path) -> str | None:
    import docx

    doc = docx.Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = " | ".join(cell.text.strip() for cell in row.cells)
            if cells.strip(" |"):
                parts.append(cells)
    return "\n".join(parts).strip()


def _extract_xlsx_text(path: Path) -> str | None:
    import openpyxl

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    try:
        parts: list[str] = []
        rows_seen = 0
        for ws in wb.worksheets:
            parts.append(f"# {ws.title}")
            for row in ws.iter_rows(values_only=True):
                if rows_seen >= _MAX_OFFICE_PREVIEW_ROWS:
                    parts.append("…")
                    break
                if any(c is not None for c in row):
                    parts.append(" | ".join("" if c is None else str(c) for c in row))
                    rows_seen += 1
        return "\n".join(parts).strip()
    finally:
        wb.close()


def _extract_pptx_text(path: Path) -> str | None:
    from pptx import Presentation

    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame and shape.text_frame.text.strip()
        ]
        if slide_lines:
            parts.append(f"**Слайд {i}:**\n" + "\n".join(slide_lines))
    return "\n\n".join(parts).strip()


_OFFICE_EXTRACTORS = {".docx": _extract_docx_text, ".xlsx": _extract_xlsx_text, ".pptx": _extract_pptx_text}


def _extract_text_preview(path: Path, content_type: str, filename: str) -> str | None:
    ext = Path(filename).suffix.lower()

    extractor = _OFFICE_EXTRACTORS.get(ext)
    if extractor is not None:
        # Реальный найденный риск security-ревью: в отличие от text/*-ветки
        # ниже (читает только первые 200 КБ), офисные экстракторы парсили
        # ВЕСЬ файл целиком синхронно в обработчике запроса, без потолка —
        # docx/pptx близко к общему лимиту 300 МБ (MAX_UPLOAD_BYTES) грузили
        # бы столько же в память на backend с бюджетом 768 МБ (CLAUDE.md).
        # xlsx уже был защищён построчным лимитом (_MAX_OFFICE_PREVIEW_ROWS)
        # — тут та же идея, но на входе, единым порогом для всех трёх.
        if path.stat().st_size > _MAX_OFFICE_PREVIEW_BYTES:
            logger.info(
                "Офисный файл %s больше %d МБ — превью пропущено",
                filename,
                _MAX_OFFICE_PREVIEW_BYTES // (1024 * 1024),
            )
            return None
        try:
            text = extractor(path)
        except Exception:
            # Битый/нестандартный файл — превью просто не будет, не должно
            # ронять саму загрузку (тот же принцип, что у PDF-превью).
            logger.exception("Не удалось извлечь текст из %s (%s)", filename, ext)
            return None
        return text[:_MAX_TEXT_PREVIEW_CHARS] if text else None

    if not (content_type.startswith("text/") or content_type == "application/json" or ext in _TEXT_PREVIEW_EXTENSIONS):
        return None
    # Читаем только "шапку" файла, не целиком — видео/архив с обманным
    # расширением не должны раздувать память ради превью, которое всё равно
    # покажет только первые символы.
    with path.open("rb") as f:
        content = f.read(_MAX_TEXT_PREVIEW_BYTES)
    text = content.decode("utf-8", errors="replace").strip()
    if not text:
        return None
    return text[:_MAX_TEXT_PREVIEW_CHARS]


def _generate_pdf_thumbnail(upload_id: uuid.UUID, pdf_bytes: bytes) -> bool:
    """Превью первой страницы PDF, кэшируется рядом с самим файлом — та же
    идея, что текстовый слой (PyMuPDF, локально, бесплатно, синхронно при
    загрузке: один рендер одной страницы дёшев, не то же самое, что
    постраничный OCR ниже). Никогда не роняет саму загрузку файла — просто
    не будет thumbnail, если PDF повреждён/нестандартный."""
    try:
        with pymupdf.open(stream=pdf_bytes, filetype="pdf") as doc:
            if doc.page_count == 0:
                return False
            pix = doc[0].get_pixmap(dpi=_MAX_THUMBNAIL_DPI)
            _thumbnail_path(upload_id).write_bytes(pix.tobytes("png"))
            return True
    except Exception:
        return False


@router.post("", response_model=UploadOut, status_code=status.HTTP_201_CREATED)
async def create_upload(
    space_id: uuid.UUID,
    file: UploadFile,
    user: User = Depends(get_current_user),
    db: AsyncSession = Depends(get_db),
) -> UploadOut:
    await ensure_space_access(db, space_id, user.id)

    content_type = file.content_type or "application/octet-stream"
    upload = Upload(
        space_id=space_id,
        author_id=user.id,
        filename=file.filename or "файл",
        content_type=content_type,
    )
    db.add(upload)
    await db.flush()

    upload_dir = Path(get_settings().upload_dir)
    upload_dir.mkdir(parents=True, exist_ok=True)
    dest = _upload_path(upload.id)

    # Пишем на диск чанками, не читаем файл целиком в память — backend
    # ограничен 768 МБ (docker-compose.yml.j2), а с потолком 300 МБ на
    # аплоад чтение целиком в bytes перед записью держало бы это в памяти
    # одним куском на весь запрос, легко несколько таких параллельно — OOM.
    total = 0
    try:
        with dest.open("wb") as out:
            while chunk := await file.read(_STREAM_CHUNK_BYTES):
                total += len(chunk)
                if total > MAX_UPLOAD_BYTES:
                    raise HTTPException(status.HTTP_400_BAD_REQUEST, "Файл больше 300 МБ")
                out.write(chunk)
    except HTTPException:
        dest.unlink(missing_ok=True)
        await db.rollback()
        raise

    # Сейф — сервер не видит plaintext (клиент шлёт уже зашифрованные
    # байты и родовое имя файла), поэтому вся обработка ниже (текстовый
    # слой PDF, превью, thumbnail, очередь OCR/расшифровки) для него
    # бессмысленна и пропускается целиком: сервер физически не может
    # прочитать то, что не может расшифровать.
    is_vault = await is_vault_space(db, space_id)
    is_pdf = content_type == "application/pdf"
    pdf_text: str | None = None
    has_thumbnail = False
    preview_text: str | None = None

    pdf_bytes: bytes | None = None
    text_layer: str | None = None
    if not is_vault:
        # Текстовый слой (если есть) — сразу, синхронно, локально (PyMuPDF,
        # без сети): дёшево независимо от размера файла. Сам ФИНАЛЬНЫЙ текст
        # (с причёсанной вёрсткой, document_reflow.py) больше не отдаётся
        # отсюда напрямую — реальная жалоба на "кашу" в вёрстке у голого
        # PyMuPDF-извлечения (особенно многоколоночные документы). Ниже он
        # идёт тем же асинхронным путём, что и сканы (пока автообработка
        # включена), только та часть, где нет сети/LLM (сам текстовый слой,
        # thumbnail, превью не-PDF), остаётся синхронной.
        pdf_bytes = dest.read_bytes() if is_pdf else None
        text_layer = extract_pdf_text(pdf_bytes) if pdf_bytes is not None else None
        has_thumbnail = _generate_pdf_thumbnail(upload.id, pdf_bytes) if pdf_bytes is not None else False
        preview_text = None if is_pdf else _extract_text_preview(dest, content_type, upload.filename)

    auto = user.auto_process_uploads and not is_vault
    pdf_ocr_queued = False
    if auto and (content_type.startswith("video/") or content_type.startswith("image/")):
        upload.transcription_status = "pending"
    elif auto and is_pdf and (text_layer or total <= AUTO_OCR_MAX_PDF_BYTES):
        # Текстовый слой — всегда (один reflow-вызов на документ, дёшево
        # независимо от размера файла). Скан без текстового слоя — раньше
        # распознавался только по кнопке «Распознать»; для небольших файлов
        # (типичный чек/страница, не многостраничная книга) дорогая
        # постраничная vision-OCR не настолько долгая/дорогая, чтобы
        # требовать лишнего клика.
        upload.transcription_status = "pending"
        pdf_ocr_queued = True
    elif is_pdf and text_layer:
        # Автообработка выключена — хотя бы сырой текстовый слой без
        # причёсывания вёрстки, лучше, чем совсем ничего; причёсанную
        # версию пользователь получит по клику «Распознать».
        pdf_text = text_layer

    await db.commit()

    if auto and content_type.startswith("video/"):
        enqueue_transcription(upload.id)
    elif auto and content_type.startswith("image/"):
        enqueue_vision(upload.id)
    elif pdf_ocr_queued:
        enqueue_pdf_ocr(upload.id)

    return UploadOut(
        id=upload.id,
        url=f"/api/uploads/{upload.id}",
        filename=upload.filename,
        content_type=content_type,
        pdf_text=pdf_text,
        pdf_ocr_queued=pdf_ocr_queued,
        preview_text=preview_text,
        has_thumbnail=has_thumbnail,
    )


# Запись прямо в заметке (кнопка-микрофон, RecordingPanel.tsx) — в отличие
# от create_upload выше, файл не приходит одним куском: MediaRecorder в
# браузере шлёт его чанками по ходу записи (защита от потери данных при
# обрыве интернета/случайном обновлении страницы — реальная просьба),
# поэтому три отдельных шага вместо одного POST. _MIN_RECORDING_BYTES —
# из note_recording.py (реальный найденный баг ручного обзора: константа
# была продублирована в двух местах, риск разъехаться при правке).


class RecordingStartIn(BaseModel):
    space_id: uuid.UUID
    # Реальный кодек, который выбрал MediaRecorder в браузере (порядок
    # предпочтений в RecordingPanel.tsx) — не фиксируем на бэкенде, он
    # разный в разных браузерах (audio/webm;codecs=opus, audio/ogg и т.п.),
    # а раздавать файл обратно (<audio src=...>) и слать в Deepgram нужно
    # с ТЕМ ЖЕ content_type, что реально записан на диск.
    content_type: str = "audio/webm"


@router.post("/recording/start", response_model=UploadOut, status_code=status.HTTP_201_CREATED)
async def start_recording(
    payload: RecordingStartIn, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> UploadOut:
    # Реальный найденный баг (ручной security-обзор): content_type раньше
    # принимался с фронта без проверки, а GET /uploads/{id} теперь отдаёт
    # файлы с Content-Disposition: inline (фикс превью PDF) — без
    # вайтлиста сюда можно было подсунуть content_type="text/html" и
    # получить stored XSS по прямой ссылке на аплоад.
    if not payload.content_type.startswith("audio/"):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "content_type записи должен быть audio/*")
    await ensure_space_access(db, payload.space_id, user.id)
    # Сама суть записи — расшифровка речи через Deepgram, которому нужен
    # реальный звук на сервере; в сейфе сервер в принципе ничего не может
    # прочитать, так что фича здесь не имеет смысла (не молча деградировать
    # до "просто плеер без расшифровки" — явно отказать).
    if await is_vault_space(db, payload.space_id):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Запись с расшифровкой недоступна в сейфе")
    upload = Upload(
        space_id=payload.space_id,
        author_id=user.id,
        filename="Запись.webm",
        content_type=payload.content_type,
        transcription_status="recording",
    )
    db.add(upload)
    await db.flush()

    upload_dir = Path(get_settings().upload_dir)
    upload_dir.mkdir(parents=True, exist_ok=True)
    _upload_path(upload.id).touch()

    await db.commit()
    return UploadOut(
        id=upload.id,
        url=f"/api/uploads/{upload.id}",
        filename=upload.filename,
        content_type=upload.content_type,
        pdf_text=None,
        pdf_ocr_queued=False,
        preview_text=None,
        has_thumbnail=False,
    )


@router.post("/recording/{upload_id}/chunk", status_code=status.HTTP_204_NO_CONTENT)
async def append_recording_chunk(
    upload_id: uuid.UUID, request: Request, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> None:
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Запись не найдена")
    await ensure_space_access(db, upload.space_id, user.id)
    if upload.transcription_status != "recording":
        raise HTTPException(status.HTTP_409_CONFLICT, "Запись уже завершена")

    body = await request.body()
    dest = _upload_path(upload_id)
    existing_size = dest.stat().st_size if dest.is_file() else 0
    if existing_size + len(body) > MAX_UPLOAD_BYTES:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Запись больше 300 МБ")
    # Порядок байт критичен: чанки MediaRecorder кроме первого не являются
    # самостоятельным валидным webm/ogg — только весь файл целиком, в том
    # порядке, в котором были записаны (RecordingPanel.tsx шлёт их строго
    # последовательно, дожидаясь ответа на предыдущий, прежде чем слать
    # следующий).
    with dest.open("ab") as f:
        f.write(body)


@router.post("/recording/{upload_id}/finish")
async def finish_recording(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> dict:
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Запись не найдена")
    await ensure_space_access(db, upload.space_id, user.id)
    if upload.transcription_status != "recording":
        raise HTTPException(status.HTTP_409_CONFLICT, "Запись уже завершена")

    dest = _upload_path(upload_id)
    size = dest.stat().st_size if dest.is_file() else 0
    if size < _MIN_RECORDING_BYTES:
        # Нажал старт и сразу стоп — почти наверняка без речи, не тратим
        # вызов Deepgram на пустую запись.
        upload.transcription_status = "failed"
        await db.commit()
        return {"status": "empty"}

    upload.transcription_status = "pending"
    await db.commit()
    enqueue_recording(upload.id)
    return {"status": "pending"}


@router.post("/{upload_id}/reprocess", status_code=status.HTTP_202_ACCEPTED)
async def reprocess_upload(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> dict:
    """Повторный запуск OCR/расшифровки — нужен файлам, загруженным до того,
    как появились vision.py/transcription.py (для них воркер никогда не
    запускался), а также если распознавание в первый раз не задалось. Для
    PDF — основной способ запустить его вручную: сканы без текстового слоя
    крупнее AUTO_OCR_MAX_PDF_BYTES или при выключенной auto_process_uploads
    автоматически не распознаются (могут быть долгими),
    только по этому явному запросу."""
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    await ensure_space_access(db, upload.space_id, user.id)
    if await is_vault_space(db, upload.space_id):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Распознавание недоступно в сейфе")

    is_pdf = upload.content_type == "application/pdf"
    if not (upload.content_type.startswith("video/") or upload.content_type.startswith("image/") or is_pdf):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Распознавание доступно только для видео, картинок и PDF")

    upload.transcription_status = "pending"
    await db.commit()

    if upload.content_type.startswith("video/"):
        enqueue_transcription(upload.id)
    elif is_pdf:
        enqueue_pdf_ocr(upload.id)
    else:
        enqueue_vision(upload.id)

    return {"status": "pending"}


@router.post("/{upload_id}/describe-now")
async def describe_upload_now(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> dict:
    """Синхронный вариант vision._process — для картинки, которую прямо
    сейчас прикрепляют к сообщению ассистенту (реальный запрос: "хочу
    скинуть фото из приложения по здоровью, чтобы ассистент разложил его
    вместе с остальным, а не пересказывать самому"). Обычный путь
    (enqueue_vision в create_upload выше) — фоновая очередь с плейсхолдером
    в content заметки; здесь плейсхолдера нет и заменять нечего, а результат
    нужен ДО отправки сообщения, не после. Не трогает transcription_status/
    transcript — если у пользователя включена auto_process_uploads, фоновая
    очередь всё равно проанализирует эту же картинку ещё раз независимо от
    этого вызова; расход на один лишний vision-запрос, не источник неверных
    данных, поэтому не координируем два пути ради разовой экономии."""
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    await ensure_space_access(db, upload.space_id, user.id)
    if not upload.content_type.startswith("image/"):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Распознавание доступно только для картинок")

    settings = get_settings()
    if not settings.llm_api_key:
        raise HTTPException(status.HTTP_503_SERVICE_UNAVAILABLE, "Распознавание временно недоступно")

    path = _upload_path(upload_id)
    if not path.is_file():
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден на диске")
    image_bytes = path.read_bytes()
    if len(image_bytes) > vision_module._MAX_IMAGE_BYTES:
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Изображение слишком большое для распознавания")

    try:
        raw = await vision_module._analyze(image_bytes, upload.content_type, settings.llm_api_key)
    except Exception:
        logger.exception("Ошибка Mistral vision при синхронном анализе %s", upload_id)
        raise HTTPException(status.HTTP_502_BAD_GATEWAY, "Не получилось распознать изображение") from None

    _, description = vision_module._split_ticket_marker(raw)
    return {"description": description.strip()}


# Вайтлист типов, которые браузер не может исполнить как активный
# контент — image/svg+xml сюда намеренно НЕ входит (может содержать
# <script>), как и text/html/application/xhtml+xml. content_type
# приходит от клиента непроверенным (create_upload), поэтому список
# должен быть закрытым, а не "всё, кроме явно плохого".
def _safe_to_inline(content_type: str) -> bool:
    if content_type == "image/svg+xml":
        return False
    return (
        content_type.startswith("image/")
        or content_type.startswith("audio/")
        or content_type.startswith("video/")
        or content_type in ("application/pdf", "text/plain")
    )


@router.put("/{upload_id}/staged", status_code=status.HTTP_204_NO_CONTENT)
async def stage_upload_replacement(
    upload_id: uuid.UUID, file: UploadFile, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> None:
    """Первая фаза смены пароля сейфа — файл перешифрован клиентом новым
    ключом и лежит здесь ВРЕМЕННОЙ копией, оригинал не тронут. Реальная
    подмена — только внутри rotate_vault_password (routers/spaces.py),
    после того как её транзакция (все заметки + новая соль/verifier)
    успешно закоммитилась. Если пользователь закроет вкладку до этого —
    орфанная .new-копия просто останется лежать, ничего не сломав (не
    более грязно, чем любой другой прерванный аплоад)."""
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    await ensure_space_access(db, upload.space_id, user.id)
    if not await is_vault_space(db, upload.space_id):
        raise HTTPException(status.HTTP_400_BAD_REQUEST, "Перешифровка доступна только для файлов в сейфе")

    dest = _staged_upload_path(upload_id)
    total = 0
    with dest.open("wb") as out:
        while chunk := await file.read(_STREAM_CHUNK_BYTES):
            total += len(chunk)
            if total > MAX_UPLOAD_BYTES:
                dest.unlink(missing_ok=True)
                raise HTTPException(status.HTTP_400_BAD_REQUEST, "Файл больше 300 МБ")
            out.write(chunk)


@router.get("/{upload_id}")
async def get_upload(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> FileResponse:
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    await ensure_space_access(db, upload.space_id, user.id)

    path = _upload_path(upload.id)
    if not path.is_file():
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    # content_disposition_type="inline" — иначе Starlette по умолчанию
    # ставит Content-Disposition: attachment из-за filename=, и браузер
    # скачивает файл вместо того, чтобы отрендерить его в <iframe>
    # (реальный найденный баг: встроенный просмотрщик PDF открывался и тут
    # же скачивал документ вместо показа). Явную кнопку "Скачать"
    # (DocumentAttachmentCard.tsx, <a download>) это не трогает — атрибут
    # download на самой ссылке форсирует сохранение независимо от
    # заголовка сервера.
    #
    # НО inline только для заведомо безопасных типов (реальный найденный
    # баг security-ревью): create_upload выше принимает content_type от
    # клиента (заголовок multipart) без всякой проверки — этот эндпоинт
    # отдаёт вообще любой загруженный файл, не только запись с диктофона,
    # так что вайтлист на audio/* (как у /recording/start) тут не подходит,
    # а безусловный inline давал stored XSS: подменил content_type на
    # text/html при загрузке — и содержимое рендерится в браузере другого
    # участника спейса как страница. Вайтлист "безопасно показать" (не
    # блэклист "опасного") — так надёжнее: неизвестный/непредвиденный тип
    # уходит в attachment, а не наоборот.
    disposition = "inline" if _safe_to_inline(upload.content_type) else "attachment"
    return FileResponse(
        path, media_type=upload.content_type, filename=upload.filename, content_disposition_type=disposition
    )


@router.delete("/{upload_id}", status_code=status.HTTP_204_NO_CONTENT)
async def delete_upload(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> None:
    """Явное удаление файла — сейчас нужен переносу заметки в сейф
    (NoteEditor.tsx): исходный НЕзашифрованный файл должен исчезнуть с
    диска сразу после того, как зашифрованная копия уже сохранена в
    сейфе, а не через 24ч грейс-период автосвипа (cleanup.py) — документ
    с перс. данными не должен лежать plaintext'ом лишние часы. Не
    проверяем, ссылается ли на файл ещё какая-то заметка (как cleanup.py
    делает для общего случая) — здесь вызывающий точно знает, что делает."""
    upload = await db.get(Upload, upload_id)
    if upload is None:
        return
    await ensure_space_access(db, upload.space_id, user.id)
    _upload_path(upload_id).unlink(missing_ok=True)
    _thumbnail_path(upload_id).unlink(missing_ok=True)
    await db.delete(upload)
    await db.commit()


@router.get("/{upload_id}/thumbnail")
async def get_upload_thumbnail(
    upload_id: uuid.UUID, user: User = Depends(get_current_user), db: AsyncSession = Depends(get_db)
) -> FileResponse:
    upload = await db.get(Upload, upload_id)
    if upload is None:
        raise HTTPException(status.HTTP_404_NOT_FOUND, "Файл не найден")
    await ensure_space_access(db, upload.space_id, user.id)

    path = _thumbnail_path(upload_id)
    if not path.is_file():
        # Файл загружен ДО того, как появилась генерация превью при
        # загрузке (_generate_pdf_thumbnail в create_upload) — досоздаём
        # лениво, по первому же запросу картинки, вместо отдельной кнопки
        # "обработать": пользователь наводит/тапает карточку — и превью
        # просто появляется само, без лишнего клика. Не бьёт по перфомансу
        # заметно — рендер одной страницы дешёвый, тот же вызов, что уже
        # синхронно делается при обычной загрузке.
        if upload.content_type != "application/pdf":
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Превью недоступно")
        src_path = _upload_path(upload_id)
        if not src_path.is_file() or not _generate_pdf_thumbnail(upload_id, src_path.read_bytes()):
            raise HTTPException(status.HTTP_404_NOT_FOUND, "Превью недоступно")
    return FileResponse(path, media_type="image/png")
